#!/usr/bin/env python3
"""
Populate an existing PPTX template with slides and speaker notes.

Usage:
  python3 tools/generate_presentation.py \
    --template "GeekTALK - Outline for Students Projects _ 4Geeks Academy.pptx" \
    --out presentation/Tapin_Correct_presentation.pptx

This script uses python-pptx. Install with:
  pip install python-pptx

It will copy the template and then add slides from a predefined outline.
"""
import argparse
from pathlib import Path
from pptx import Presentation


OUT_DIR = Path("presentation")
OUT_DIR.mkdir(exist_ok=True)

SLIDES = [
    ("Title", "Tapin_Correct — React + Flask Marketplace"),
    (
        "Elevator Pitch",
        "Full‑stack template: React (Vite) front end + Flask API. Includes auth, listings, email flows, and ready-for-deploy scripts.",
    ),
    (
        "Architecture",
        "Frontend: src/front (Vite + React). Backend: src/backend (Flask + SQLAlchemy). Build -> dist served by Flask.",
    ),
    (
        "Backend Features",
        "JWT auth, listings endpoints, Alembic migrations, seed scripts, email/password reset flows.",
    ),
    (
        "Frontend Features",
        "Component-based SPA, map integration (Leaflet), Playwright E2E tests, responsive UI.",
    ),
    (
        "Developer Experience",
        "Pipenv/venv for backend, Node + Vite for frontend, .env-driven configuration and root package.json scripts.",
    ),
    (
        "Testing & Quality",
        "Pytest for backend, Playwright for E2E, pre-commit hooks (Black, detect-secrets).",
    ),
    (
        "Deployment",
        "Build frontend (npm run build) -> dist, set env vars on host (Render/Heroku), serve via WSGI.",
    ),
    (
        "Demo Plan",
        "Start backend, start frontend, demo register/login, create listing, view map, run basic tests.",
    ),
    (
        "Production Notes",
        "Use managed DB (Postgres), strong secrets, HTTPS, CORS and rate-limiting.",
    ),
    ("Roadmap", "Payments, advanced search, recommendations, CI/CD, expanded tests."),
    (
        "Resources",
        "README.md, src/backend/README.md, 4Geeks starter docs. Contact for follow-up.",
    ),
]


def add_text_slide(prs, title, content):
    # Use the 'Title and Content' layout (index 1) from the template
    title_and_content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for line in content.split("\n"):
        p = body.add_paragraph()
        p.text = line
    return slide


def populate_presentation(template_path: Path, out_path: Path):
    if not template_path.exists():
        raise SystemExit(f"Template not found: {template_path}")
    prs = Presentation(str(template_path))

    # Clear existing slides except the title slide
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    for title, content in SLIDES:
        slide = add_text_slide(prs, title, content)
        # Add speaker notes
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"Speaker notes — {title}: {content}"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Created presentation: {out_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", required=True)
    parser.add_argument(
        "--out", default=str(OUT_DIR / "Tapin_Correct_presentation.pptx")
    )
    args = parser.parse_args()
    template = Path(args.template)
    out = Path(args.out)
    populate_presentation(template, out)


if __name__ == "__main__":
    main()
