#!/usr/bin/env python3
"""
Inspect a PPTX file and list its slide layouts and placeholders.

Usage:
  python3 tools/inspect_pptx.py "path/to/presentation.pptx"
"""
import sys
from pptx import Presentation


def inspect_presentation(path):
    """List slide layouts and the placeholders in each."""
    prs = Presentation(path)
    print(f"Inspecting template: {path}")
    print("-" * 30)
    print(f"Found {len(prs.slides)} slides and {len(prs.slide_layouts)} layouts.")
    print("-" * 30)

    for i, layout in enumerate(prs.slide_layouts):
        print(f"\nLayout {i}: {layout.name}")
        try:
            for j, ph in enumerate(layout.placeholders):
                print(
                    f"  - Placeholder {j}: {ph.name} ({ph.placeholder_format.idx}), type {ph.placeholder_format.type}"
                )
        except Exception as e:
            print(f"  - Error inspecting placeholders for layout {i}: {e}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 tools/inspect_pptx.py <path_to_pptx>")
        sys.exit(1)
    inspect_presentation(sys.argv[1])
